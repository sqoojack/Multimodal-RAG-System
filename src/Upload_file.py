# Upload_file.py
from utils import io, pytesseract, base64, requests, re, json, ffmpeg, tempfile, os, whisper
from langchain_core.documents import Document
from pptx import Presentation
from docx import Document as DocxDocument
from PyPDF2 import PdfReader
from PIL import Image, ImageFile
from config import ollama_url

# Generate image description using multimodal LLM (e.g. LLaVA)
def get_image_caption(image_bytes, ollama_url, img_model):
    # Convert image to base64 format for Ollama API
    b64_image = base64.b64encode(image_bytes).decode()

    # Create POST request payload
    payload = {
        "model": img_model,
        "prompt": "Please describe the content of this image to help visually impaired individuals understand it.",
        "images": [b64_image]
    }

    # Send request to Ollama API
    response = requests.post(f"{ollama_url}/api/generate", json=payload)

    # Combine response lines into a single string
    result = ""
    for line in response.iter_lines():
        if line:
            result += line.decode("utf-8")

    responses = []

    # Extract JSON objects using regex
    json_objects = re.findall(r'\{.*?"response":.*?\}', result)

    # Parse JSON objects and extract response field
    for obj in json_objects:
        try:
            data = json.loads(obj)
            if "response" in data:
                responses.append(data["response"])
        except json.JSONDecodeError:
            continue  # Ignore invalid JSON segments

    # Return combined image description
    return "".join(responses).strip()

# Load uploaded file and convert to Langchain Document format
def load_documents_from_upload(uploaded_file, img_model):

    # Get filename and content
    filename = uploaded_file.name.lower()
    content = uploaded_file.read()

    # === Handle PDF files ===
    if filename.endswith(".pdf"):
        reader = PdfReader(io.BytesIO(content))  # Read binary content
        docs = []
        for idx, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():  # Only add pages with text
                docs.append(Document(
                    page_content=text.strip(),
                    metadata={
                        "source": uploaded_file.name,
                        "type": "pdf",
                        "page": idx + 1  # 1-based page index
                    }
                ))
        return docs

    # === Handle image files (jpg, jpeg, png) ===
    elif filename.endswith((".jpg", ".jpeg", ".png")):
        ImageFile.LOAD_TRUNCATED_IMAGES = True  # Allow loading truncated images
        image = Image.open(io.BytesIO(content)).convert("RGB")  # Convert to RGB mode

        # Perform OCR using Tesseract (Traditional Chinese + English)
        ocr_text = pytesseract.image_to_string(image, lang="chi_tra+eng")
        ocr_lines = [l.strip() for l in ocr_text.splitlines() if l.strip()]
        ocr_paragraph = "\n".join(ocr_lines)

        # Get image caption from LLM
        caption = get_image_caption(content, ollama_url, img_model) or "(No response)"

        # Combine caption and OCR text
        full_text = f"""【Image Description】\n{caption}\n\n【OCR Text】\n{ocr_paragraph}"""

        return [
            Document(
                page_content=full_text.strip(),
                metadata={
                    "source": uploaded_file.name,
                    "type": "image+ocr+llava"
                }
            )
        ]
    
    # === Handle DOCX files ===
    elif filename.endswith(".docx"):
        text = ""
        doc = DocxDocument(io.BytesIO(content))
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text.strip() + "\n"
        return [Document(page_content=text.strip(),
                    metadata={"source": uploaded_file.name, "type": "doc"})]
    
    # === Handle legacy DOC files ===
    elif filename.endswith(".doc"):   
        # 1. Write to temporary file
        with tempfile.TemporaryDirectory() as tmp:
            path_doc = os.path.join(tmp, "input.doc")
            with open(path_doc, "wb") as f:
                f.write(content)

            # 2. Read text via Word COM automation
            word = client.Dispatch("Word.Application")
            word.Visible = False
            doc  = word.Documents.Open(path_doc, ReadOnly=True)
            text = doc.Content.Text
            doc.Close(False)
            word.Quit()

        # 3. Return LangChain Document
        return [
            Document(
                page_content=text.strip(),
                metadata={"source": uploaded_file.name, "type": "doc"})]
    
    # === Handle PPTX and PPT files ===
    elif filename.endswith((".pptx", ".ppt")):
        prs = Presentation(io.BytesIO(content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = shape.text.strip()
                    if txt:
                        text += txt + "\n"
        return [Document(page_content=text.strip(),
                        metadata={"source": uploaded_file.name, "type": "ppt"})]
    
    # === Handle MP3 files ===
    elif filename.endswith(".mp3"):
        with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as tmp:
            tmp.write(content)
            tmp_path = tmp.name
        whisper_model = whisper.load_model("base")
        result = whisper_model.transcribe(tmp_path)    # Transcribe audio using Whisper
        transcript = result["text"].strip()
        
        return [Document(page_content=transcript, metadata={"source": uploaded_file.name, "type": "audio"})]

    # === Handle MP4 files ===
    elif filename.endswith(".mp4"):
        # 1. Create a temporary MP4 file
        with tempfile.NamedTemporaryFile(suffix=".mp4", delete=False) as tmp:
            tmp.write(content) # Write uploaded bytes to file
            tmp_path = tmp.name
        
        try:
            # 2. Transcribe directly using Whisper
            whisper_model = whisper.load_model("base")
            result = whisper_model.transcribe(tmp_path)
            transcript = result["text"].strip()
            
            return [Document(page_content=transcript, metadata={"source": uploaded_file.name, "type": "audio"})]
        
        except Exception as e:
            print(f"Transcribe error: {e}")
            raise e
        
        finally:
            # 3. Remove temporary file
            if os.path.exists(tmp_path):
                os.remove(tmp_path)

    # === Handle source code files (.py, .c, .cpp) ===
    elif filename.endswith((".py", ".c", ".cpp")):
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("big5", errors="ignore")

        # Keep original spacing and line breaks for source code
        return [
            Document(
                page_content=text,
                metadata={
                    "source": uploaded_file.name,
                    "type": "code"
                }
            )
        ]

    # === Handle plain text files (.txt) ===
    else:
        try:
            text = content.decode("utf-8")         # Decode using UTF-8
        except UnicodeDecodeError:
            text = content.decode("big5", errors="ignore")      # Fallback to Big5

        # Remove empty lines and reconstruct text
        full_text = "\n".join([l.strip() for l in text.splitlines() if l.strip()])

        return [
            Document(
                page_content=full_text,
                metadata={
                    "source": uploaded_file.name,
                    "type": "txt"
                }
            )
        ]


# Process multiple files and aggregate LangChain documents
def process_uploaded_files(uploaded_files, img_model):
    docs = []
    for f in uploaded_files or []:
        docs.extend(load_documents_from_upload(f, img_model))   # Append documents sequentially
    return docs